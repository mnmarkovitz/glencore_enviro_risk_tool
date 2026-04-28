"""
Build a 7-slide PPTX deck for the graduate-level class presentation.
Speaker = 2 people, 5 min each. Geopolitics-of-energy lens.
Output: Glencore_Env_Risk_Tool_Deck.pptx (uploadable to Google Slides).

Slides:
 1. Title
 2. Scope of Work (per Glencore × NYU CGA Practicum brief)
 3. Why this matters now — energy transition × supply concentration
 4. Critical-mineral × CAHRA overlap (evidence)
 5. The tool — overview (live screenshot)
 6. Methodology — formulas + data sources
 7. Geopolitics deep dive — DRC × Cobalt × Mining (case study)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

ROOT = Path(__file__).parent.parent
ASSETS = ROOT / "docs" / "slide_assets"
OUT = ROOT / "Glencore_Env_Risk_Tool_Deck.pptx"

# --- Glencore palette ---
TEAL = RGBColor(0x00, 0xA9, 0xA5)
DEEP = RGBColor(0x00, 0x5F, 0x73)
LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
RED = RGBColor(0xE5, 0x39, 0x35)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
AMBER = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=18, color=TEXT_DARK, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
              font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=TEXT_DARK,
                 line_spacing=1.18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = line_spacing
        bold = False
        text = item
        if isinstance(item, tuple):
            text, bold = item
        run = para.add_run()
        run.text = "• " + text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"
    return tb


def add_header(slide, title, subtitle=None):
    """Top accent bar + title."""
    add_rect(slide, 0, 0, SW, Inches(0.15), TEAL)
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
              title, size=26, color=DEEP, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                  subtitle, size=14, color=TEXT_GREY, italic=True)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
              "Glencore × NYU SPS Center for Global Affairs Consulting Practicum",
              size=9, color=TEXT_GREY, italic=True)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35),
              f"{page_num} / {total}", size=9, color=TEXT_GREY,
              align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


# ========================================================================
# Slide 1 — Title
# ========================================================================
s = prs.slides.add_slide(BLANK)
# Full teal background
add_rect(s, 0, 0, SW, SH, DEEP)
add_rect(s, 0, Inches(2.5), SW, Inches(2.5), TEAL)
add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.6),
          "ENVIRONMENTAL RISK IDENTIFICATION & ASSESSMENT TOOL",
          size=16, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.6),
          "For Glencore's Group Responsible Sourcing team",
          size=40, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
          "A geopolitics-of-energy lens on critical-mineral supply chains",
          size=18, color=LIGHT, italic=True)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
          "Glencore × NYU SPS Center for Global Affairs Consulting Practicum",
          size=14, color=LIGHT, italic=True)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.3),
          "Spring 2026",
          size=11, color=LIGHT)
add_speaker_notes(s,
    "[Speaker 1, 30 sec] Welcome. We're presenting an environmental risk "
    "identification tool we built for the Glencore Responsible Sourcing team. "
    "Our framing is the geopolitics of the energy transition — the minerals "
    "we need for batteries and grids are concentrated in jurisdictions where "
    "environmental due diligence is hardest. We'll walk through the scope, "
    "the geopolitical context, the tool itself, a case study, and how it "
    "fits Glencore's existing workflow.")


# ========================================================================
# Slide 2 — Scope of Work
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Scope of Work",
            "Glencore × NYU SPS Consulting Practicum — Environmental Due Diligence Framework Development")

# Top — challenge framing
add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.95), LIGHT)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.35),
          "The challenge", size=14, color=DEEP, bold=True)
add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.55),
          "Integrate environmental risk into Glencore's existing human-rights-focused due diligence "
          "framework. 200,000+ goods & services suppliers and 1,000+ metals & minerals suppliers; "
          "small assessment team that is not a subject-matter expert across every risk type. "
          "→ Scalable, risk-based approach essential.",
          size=11, color=TEXT_DARK)

# Two columns: project objective + key research questions
add_text(s, Inches(0.5), Inches(2.65), Inches(6.0), Inches(0.4),
          "Project objective", size=14, color=DEEP, bold=True)
add_bullets(s, Inches(0.5), Inches(3.05), Inches(6.0), Inches(2.0), [
    ("Develop a practical, risk-based environmental due-diligence "
     "guideline and process workflow", True),
    "Build on Glencore's saliency-based risk framework",
    "Account for very limited assessor resources",
    "Inform compliance with EU CSRD, CSDDD, EU Battery Regulation, and the "
    "Consolidated Mining Standard Initiative",
], size=11)

add_text(s, Inches(7.0), Inches(2.65), Inches(5.8), Inches(0.4),
          "Key research questions", size=14, color=DEEP, bold=True)
add_bullets(s, Inches(7.0), Inches(3.05), Inches(5.8), Inches(2.0), [
    "How do we assess severity of environmental risk?",
    "How do we prioritize among competing concerns?",
    ("How do we translate risk identification into actionable due "
     "diligence with a small team?", True),
    "How does sector + jurisdiction influence prioritization "
    "(water scarcity, protected areas, hazardous materials)?",
], size=11)

# Bottom strip: deliverables + frameworks
add_text(s, Inches(0.5), Inches(5.15), Inches(6.0), Inches(0.4),
          "Deliverables", size=14, color=DEEP, bold=True)
add_bullets(s, Inches(0.5), Inches(5.55), Inches(6.0), Inches(1.0), [
    "1. Environmental Due Diligence Guideline",
    "2. End-to-end process workflow (visual)",
    "3. This presentation + tool prototype",
], size=11)

add_text(s, Inches(7.0), Inches(5.15), Inches(5.8), Inches(0.4),
          "Anchor frameworks", size=14, color=DEEP, bold=True)
add_bullets(s, Inches(7.0), Inches(5.55), Inches(5.8), Inches(1.0), [
    "OECD Handbook on Environmental DD in Mineral Supply Chains (2023)",
    "RMI Supply Chain DD Plus Module (April 2025)",
    "UNDP HRDD and the Environment: A Practical Tool for Business",
    "Glencore's saliency-based risk framework + SCDD M&M Procedure",
], size=11)

# Cite footer
add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
          "Source: Glencore × NYU Consulting Practicum project brief, NYU SPS Center for Global Affairs, 2026.",
          size=9, color=TEXT_GREY, italic=True)
add_footer(s, 2, 7)
add_speaker_notes(s,
    "[Speaker 1, 1 min] Glencore — a global producer, processor, and marketer "
    "of transition-enabling metals and minerals — partnered with NYU's Center "
    "for Global Affairs to address a specific challenge: their existing "
    "due-diligence framework is human-rights-focused, but emerging EU "
    "regulation (CSRD, CSDDD, the Battery Regulation, the Consolidated Mining "
    "Standard) demands integrated environmental risk assessment. The "
    "operational constraint is acute: Glencore screens over 200,000 goods "
    "and services suppliers and 1,000+ metals and minerals suppliers with "
    "a very small team that cannot be subject-matter experts across every "
    "risk type. So scalability and a risk-based approach are essential. "
    "We anchored our work in the OECD Environmental DD Handbook, the RMI "
    "Plus Module, the UNDP Practical Tool for Business, and Glencore's own "
    "saliency framework. The deliverables are the guideline document, the "
    "end-to-end workflow, and what we'll show you today: a working tool "
    "prototype that operationalizes Tier 1. I'll now hand to [Speaker 2] "
    "for the geopolitical framing.")


# ========================================================================
# Slide 3 — Geopolitics framing
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Why this matters now",
            "The energy transition is mineral-intensive — and supply is concentrated")
# Big stat callouts
def stat(slide, x, y, big, label, color=TEAL):
    add_rect(slide, x, y, Inches(3.1), Inches(2.0), color)
    add_text(slide, x, y + Inches(0.2), Inches(3.1), Inches(0.9), big,
              size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.15), Inches(3.1), Inches(0.8), label,
              size=11, color=WHITE, align=PP_ALIGN.CENTER)

stat(s, Inches(0.5), Inches(1.55), "4–6×",
      "IEA: critical-mineral demand growth by 2040 under net-zero scenarios", DEEP)
stat(s, Inches(3.7), Inches(1.55), "74%",
      "Cobalt — DRC share of 2023 global production (USGS MCS 2024)", TEAL)
stat(s, Inches(6.9), Inches(1.55), "48%",
      "Nickel — Indonesia share of 2023 global production", TEAL)
stat(s, Inches(10.1), Inches(1.55), "72%",
      "Platinum — South Africa share of 2023 global production", DEEP)

# Text block
add_text(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
          "Two trends collide", size=16, color=DEEP, bold=True)
add_bullets(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.0), [
    ("(a) The energy transition runs through critical minerals.", True),
    "Cobalt + nickel for batteries · copper for grid · REEs for wind turbines · "
    "platinum for fuel cells. The IEA estimates 4–6× demand growth by 2040.",
    ("(b) Production is concentrated — and concentration overlaps with "
     "high-risk jurisdictions.", True),
    "The US 2022 Critical Minerals List, the EU Critical Raw Materials Act (2024), "
    "and the IRA all flag the same supply-concentration risk. For an operator "
    "like Glencore, that is the operational reality, not an abstraction.",
], size=12)

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
          "Sources: IEA Critical Minerals Outlook 2024 (iea.org); "
          "USGS Mineral Commodity Summaries 2024 (pubs.usgs.gov/periodicals/mcs2024); "
          "USGS 2022 Critical Minerals List; EU Critical Raw Materials Act 2024.",
          size=9, color=TEXT_GREY, italic=True)
add_footer(s, 3, 7)
add_speaker_notes(s,
    "[Speaker 2, 1 min 30] The energy transition isn't just a story about "
    "renewable electricity — it's a mineral story. The IEA forecasts 4 to 6 "
    "times demand growth for cobalt, nickel, copper, and rare earths by 2040 "
    "under a net-zero scenario. But supply is highly concentrated: 74 percent "
    "of cobalt comes from the DRC, almost half of nickel from Indonesia, "
    "72 percent of platinum from South Africa. And here's the critical "
    "geopolitical insight — the US 2022 Critical Minerals List, the EU's "
    "Critical Raw Materials Act, and the Inflation Reduction Act are all "
    "responding to the same problem: every Western government is now treating "
    "supply-chain concentration as a national-security risk. For an operator "
    "like Glencore, that's not a thought experiment — it's where they actually "
    "buy and sell every day.")


# ========================================================================
# Slide 4 — Supply concentration meets CAHRA
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The CAHRA × Critical Minerals overlap",
            "Where the energy transition's supply chain meets Glencore's high-risk list")

# Build a small table of overlap evidence
table_data = [
    ("Commodity", "Top producer", "Share", "CAHRA", "Critical?"),
    ("Cobalt", "DRC", "74%", "All regions", "Yes"),
    ("Nickel", "Indonesia", "48%", "All regions", "Yes"),
    ("Platinum", "South Africa", "72%", "—", "Yes"),
    ("Manganese alloys", "South Africa / Gabon", "36% + 23%", "—", "Yes"),
    ("Ferrochrome", "South Africa", "44%", "—", "Yes"),
    ("Aluminum", "China (smelting)", "59%", "Xinjiang", "Yes"),
    ("Zinc", "China", "33%", "Xinjiang", "Yes"),
    ("Copper", "Chile / Peru / DRC", "23 + 10 + 13%", "DRC: All regions", "—"),
]
# Manual table
left = Inches(0.5); top = Inches(1.5); col_w = [Inches(2.4), Inches(3.0),
    Inches(2.0), Inches(2.8), Inches(1.8)]
total_w = sum(c for c in col_w)
row_h = Inches(0.42)
for r_i, row in enumerate(table_data):
    is_header = r_i == 0
    x = left
    for c_i, val in enumerate(row):
        cell = add_rect(s, x, top + row_h * r_i, col_w[c_i], row_h,
                         DEEP if is_header else (LIGHT if r_i % 2 == 1 else WHITE),
                         line=DEEP if is_header else TEXT_GREY)
        add_text(s, x, top + row_h * r_i, col_w[c_i], row_h, val,
                  size=12, color=WHITE if is_header else TEXT_DARK,
                  bold=is_header, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[c_i]

# Insight callout
add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.0), TEAL)
add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.7),
          "Of the 8 USGS-Critical Minerals in Glencore's portfolio, every single one has a top-3 "
          "producer that touches the CAHRA list. The energy transition cannot be decoupled from "
          "high-risk jurisdictions — managing environmental risk in those countries is non-optional.",
          size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
          "Sources: USGS MCS 2024; USGS 2022 Critical Minerals List "
          "(usgs.gov/news/national-news-release/us-geological-survey-releases-2022-list-critical-minerals); "
          "Glencore CAHRA List 2025 (updated 27.02.2025).",
          size=9, color=TEXT_GREY, italic=True)
add_footer(s, 4, 7)
add_speaker_notes(s,
    "[Speaker 2, 1 min 30] If we map Glencore's CAHRA list against the US "
    "2022 Critical Minerals List, the overlap is striking. Cobalt — 74 percent "
    "DRC, all of which is CAHRA-flagged. Nickel — 48 percent Indonesia, "
    "all CAHRA. Platinum — 72 percent South Africa. Aluminum smelting — "
    "59 percent China, with Xinjiang flagged. The point isn't that any "
    "specific country is bad; it's that the geographic concentration of "
    "energy-transition minerals systematically pushes Western buyers into "
    "jurisdictions where environmental governance is weak or contested. "
    "That's the operational gap our tool addresses. Now I'll hand to "
    "[Speaker 1] for the tool itself.")


# ========================================================================
# Slide 5 — The tool overview (screenshot)
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The tool",
            "Pivot-style filter → ranked risks · Likelihood × Severity heatmap · drill-down audit trail")

# Add screenshot (left half)
img = ASSETS / "01_dashboard.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.4), Inches(1.3), height=Inches(5.5))

# Right side commentary
add_text(s, Inches(8.0), Inches(1.4), Inches(4.9), Inches(0.4),
          "What it does", size=16, color=DEEP, bold=True)
add_bullets(s, Inches(8.0), Inches(1.85), Inches(4.9), Inches(4.5), [
    "Filters: commodity, country, process, risk type — top producers surfaced first",
    "8 priority risks (15 total) scored as Likelihood × Severity = Overall (1–25)",
    "Bucketed Low / Moderate / High / Critical with consistent palette across "
    "the table, heatmap, choropleth, and Excel companion",
    "Drill-down panel exposes raw indicator + source URL for every score → "
    "auditable for CSRD / CSDDD",
    "8 tabs: Dashboard · Map · Comparative Analysis · Risk Library · "
    "SCDD Tiers · User Guide · Methodology · Data Sources",
], size=11)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.3),
          "Live tool: streamlit.app  ·  Repo: github.com/mnmarkovitz/glencore_enviro_risk_tool",
          size=9, color=TEXT_GREY, italic=True)
add_footer(s, 5, 7)
add_speaker_notes(s,
    "[Speaker 1, 1 min] This is the tool. An analyst opens it, picks a "
    "commodity in the sidebar — let's say cobalt — and the country "
    "dropdown automatically prioritizes the top producers from the USGS "
    "Mineral Commodity Summaries. They pick a process — mining, refining, "
    "smelting, recycling, or marketing — and they get back a ranked table "
    "of every environmental risk that applies, plus a 5-by-5 Likelihood-"
    "by-Severity heatmap. Click any row and you see the raw indicator value "
    "with its source URL, so an audit can trace every number back to its "
    "public data point.")


# ========================================================================
# Slide 6 — Methodology
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Methodology",
            "Reproducible scoring from public datasets · No analyst gut-feel")

# Formulas
add_rect(s, Inches(0.5), Inches(1.5), Inches(8.5), Inches(2.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.65), Inches(8.1), Inches(0.4),
          "Two formulas", size=14, color=DEEP, bold=True)
add_text(s, Inches(0.7), Inches(2.1), Inches(8.1), Inches(0.5),
          "Likelihood = 0.4 × Process Intrinsic + 0.6 × Country Hazard    →   1–5",
          size=14, color=TEXT_DARK, font="Consolas")
add_text(s, Inches(0.7), Inches(2.55), Inches(8.1), Inches(0.5),
          "Severity   = 0.5 × Ecological Sensitivity + 0.5 × Regulatory Strict.   →   1–5",
          size=14, color=TEXT_DARK, font="Consolas")
add_text(s, Inches(0.7), Inches(3.05), Inches(8.1), Inches(0.5),
          "Overall    = Likelihood × Severity              →   1–25",
          size=14, color=TEXT_DARK, font="Consolas")
add_text(s, Inches(0.7), Inches(3.55), Inches(8.1), Inches(0.4),
          "Buckets: 1–4 Low · 5–9 Moderate · 10–14 High · 15–25 Critical",
          size=11, color=TEXT_GREY, italic=True)

# Bucket strip
def pill(slide, x, y, w, h, label, color):
    add_rect(slide, x, y, w, h, color)
    add_text(slide, x, y, w, h, label, size=11, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

pill(s, Inches(9.3), Inches(1.5), Inches(0.85), Inches(0.35), "Low", GREEN)
pill(s, Inches(10.2), Inches(1.5), Inches(0.85), Inches(0.35), "Moderate", AMBER)
pill(s, Inches(11.1), Inches(1.5), Inches(0.85), Inches(0.35), "High", ORANGE)
pill(s, Inches(12.0), Inches(1.5), Inches(0.85), Inches(0.35), "Critical", RED)

# Why these weights
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
          "Why country weight > process weight (0.6 vs 0.4)?",
          size=14, color=DEEP, bold=True)
add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.0),
          "Process determines whether a risk is possible at all (no tailings in trading), but country "
          "context determines whether a capable process actually causes harm. Two copper mines — Chile's "
          "Atacama and Zambia's Copperbelt — have the same Process Intrinsic score for water depletion "
          "(5), but very different realized risk: Aqueduct rates Atacama 5 (Extremely High) vs Zambia 2.",
          size=12, color=TEXT_DARK)

# Data sources strip
add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
          "Public datasets feeding the scores",
          size=14, color=DEEP, bold=True)
add_text(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.5),
          "WRI Aqueduct · Yale EPI 2024 · WHO PM2.5 DB · IUCN Red List · WDPA · Global Forest Watch · "
          "Global Tailings Portal · USGS MRDS · USGS Critical Minerals Atlas · Global Energy Monitor · "
          "ISRIC SoilGrids · World Bank WGI · INFORM Risk · Glencore CAHRA List 2025",
          size=10, color=TEXT_GREY, italic=True)

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
          "All weights and bucket thresholds editable in scoring_weights.csv. "
          "Full methodology: github.com/mnmarkovitz/glencore_enviro_risk_tool/blob/main/docs/METHODOLOGY.md",
          size=9, color=TEXT_GREY, italic=True)
add_footer(s, 6, 7)
add_speaker_notes(s,
    "[Speaker 1, 1 min] Two formulas. Likelihood is 40 percent process "
    "intensity — does the process even cause this risk — and 60 percent "
    "country hazard, drawn from a specific public dataset for each risk: "
    "Aqueduct for water, IUCN for species, GFW for deforestation, "
    "SoilGrids for soil, and so on. Severity is half ecological sensitivity, "
    "half regulatory strictness — because the tool measures penalty exposure, "
    "not pure ecological damage. Overall is L times S, bucketed into Low, "
    "Moderate, High, and Critical. Importantly, every single weight is in an "
    "editable CSV — Glencore can change the balance whenever they want without "
    "touching code. I'll hand to [Speaker 2] for the case study.")


# ========================================================================
# Slide 7 — Geopolitics deep dive — DRC × Cobalt × Mining
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Geopolitics deep dive — DRC × Cobalt × Mining",
            "Why this combination is the canonical stress-test for the energy transition")

# Top row: 4 fact tiles framing the stakes
def fact_tile(slide, x, y, w, h, num, label, color):
    add_rect(slide, x, y, w, h, color)
    add_text(slide, x, y + Inches(0.1), w, Inches(0.8), num,
              size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.95), w, Inches(0.8), label,
              size=10, color=WHITE, align=PP_ALIGN.CENTER)

fact_tile(s, Inches(0.5), Inches(1.4), Inches(3.05), Inches(1.65),
           "74%", "DRC's share of global cobalt mine production (USGS MCS 2024)", DEEP)
fact_tile(s, Inches(3.65), Inches(1.4), Inches(3.05), Inches(1.65),
           "All regions", "Glencore CAHRA flag — every Congolese province is high-risk", TEAL)
fact_tile(s, Inches(6.8), Inches(1.4), Inches(3.05), Inches(1.65),
           "~2 of 5", "of every EV battery's cathode chemistry depends on cobalt", DEEP)
fact_tile(s, Inches(9.95), Inches(1.4), Inches(3.05), Inches(1.65),
           "Tier 2",
           "auto-trigger: any DRC supplier escalates immediately under SCDD M&M procedure", TEAL)

# Geopolitics narrative — left column
add_text(s, Inches(0.5), Inches(3.2), Inches(6.0), Inches(0.4),
          "Why DRC × cobalt is structurally exposed",
          size=14, color=DEEP, bold=True)
add_bullets(s, Inches(0.5), Inches(3.6), Inches(6.0), Inches(3.2), [
    ("Energy transition concentration: cobalt for EV batteries, grid storage. "
     "IEA forecasts demand 6× by 2040.", True),
    ("Geographic concentration: DRC produces ~74%; refining ~70% China-controlled. "
     "Western buyers face cross-strait + cross-border policy exposure.", True),
    "EU Battery Regulation (EU 2023/1542) and the U.S. Inflation Reduction Act both treat "
    "cobalt traceability and DD as compliance gating items for market access.",
    "Artisanal & small-scale mining (ASM) supplies ~15-30% of DRC cobalt — child labor and "
    "fatality risk make environmental + human-rights DD inseparable.",
    "Acidic tropical Oxisols → highest soil-pollution vulnerability in the dataset; "
    "spilled heavy metals mobilize fast (ISRIC SoilGrids vulnerability ≈ 4.3 / 5).",
], size=11)

# Tool screenshot — right column
img = ASSETS / "01_dashboard.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(6.85), Inches(3.2), height=Inches(3.4))
add_text(s, Inches(6.85), Inches(6.55), Inches(6.0), Inches(0.3),
          "Live tool view: DRC × Cobalt × Mining → 3 Critical risks",
          size=10, color=TEXT_GREY, italic=True)

# Bottom — what the tool delivers operationally
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
          "What the analyst gets in 10 minutes: ranked risks (Tailings · Soil pollution · "
          "Biodiversity → all Critical), KPIs for the SAQ, supplier-type categories, every "
          "score auditable to a public dataset URL.",
          size=10, color=TEXT_DARK)

# Sources strip
add_text(s, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.3),
          "Sources: USGS MCS 2024 · Glencore CAHRA List 2025 · IEA Critical Minerals Outlook 2024 · "
          "ISRIC SoilGrids 2.0 · IUCN Red List · Global Forest Watch · Global Tailings Portal · "
          "EU Battery Regulation 2023/1542 · U.S. IRA (Public Law 117–169).",
          size=8, color=TEXT_GREY, italic=True)
add_footer(s, 7, 7)
add_speaker_notes(s,
    "[Speaker 2, 2 min — geopolitics deep dive] DRC × cobalt × mining is "
    "the single best case study to show why this tool matters in the "
    "geopolitics of the energy transition. Three layers stack here. "
    "First, energy-transition concentration: every electric vehicle "
    "battery cathode chemistry except LFP relies on cobalt, and the IEA "
    "forecasts 6× demand growth by 2040. Second, geographic concentration: "
    "74 percent of mine production is in the DRC, and 70 percent of "
    "refining is in China — so a Western buyer is dependent on a single "
    "high-risk source country plus a single processing chokepoint that "
    "is itself subject to escalating U.S.–China trade restrictions. "
    "Third, regulatory pressure: the EU Battery Regulation that came "
    "into force in 2023 and the U.S. Inflation Reduction Act both treat "
    "cobalt traceability and environmental due diligence as gating items "
    "for market access. So this isn't optional — Glencore needs auditable "
    "Tier-1 evidence per supplier. The tool surfaces three Critical risks "
    "for any DRC cobalt mine: tailings, soil pollution — because DRC's "
    "tropical Oxisols mobilize heavy metals fast — and biodiversity loss. "
    "Every score links back to a public dataset, which is exactly what an "
    "EU CSDDD or CSRD audit will demand. That collapses a multi-day "
    "OSDR exercise into about ten minutes per supplier — and that's the "
    "scalability gap Glencore's small assessment team needs filled.")


prs.save(OUT)
print(f"✓ Wrote {OUT.name}  ({OUT.stat().st_size / 1024:.0f} KB)")
print(f"  Location: {OUT}")
